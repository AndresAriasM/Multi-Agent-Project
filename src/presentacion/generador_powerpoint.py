"""
Generador de presentaciones PowerPoint ULTRA DETALLADO
Ubicación: src/presentacion/generador_powerpoint.py
Incluye TODAS las instituciones, programas y datos disponibles
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, List
import pandas as pd
import os
from pathlib import Path


class GeneradorPowerPoint:
    """Genera presentaciones PowerPoint ULTRA DETALLADAS con datos enriquecidos"""
    
    def __init__(self, datos: Dict, resultados_agentes: Dict, graficas_dir: str = 'output'):
        self.datos = datos
        self.resultados = resultados_agentes
        self.graficas_dir = graficas_dir
        self.datos_enriquecidos = datos.get('datos_enriquecidos', {})
        
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        
        # Colores corporativos
        self.color_primario = RGBColor(31, 78, 121)
        self.color_secundario = RGBColor(79, 129, 189)
        self.color_acento = RGBColor(192, 0, 0)
        self.color_texto = RGBColor(0, 0, 0)
        self.color_blanco = RGBColor(255, 255, 255)
        self.color_exito = RGBColor(46, 204, 113)
        self.color_advertencia = RGBColor(230, 126, 34)
    
    def crear_presentacion(self, filepath: str) -> None:
        """Crea la presentación completa con datos ultra detallados"""
        print("📊 Generando presentación PowerPoint ULTRA detallada...")
        
        try:
            # Diapositivas básicas
            self._agregar_portada()
            self._agregar_tabla_contenidos()
            
            # Resumen ejecutivo
            self._agregar_resumen_ejecutivo()
            
            # Datos enriquecidos - SECCIÓN EXTENDIDA
            if self.datos_enriquecidos:
                self._agregar_contexto_mercado()
                
                # SECCIÓN 1: PROGRAMAS DETALLADO
                self._agregar_programas_equivalentes_completo()
                self._agregar_denominaciones_analisis()
                
                # SECCIÓN 2: INSTITUCIONES DETALLADO
                self._agregar_instituciones_completo()
                self._agregar_universidades_oferentes()
                self._agregar_tecnologicas_oferentes()
                self._agregar_distribucion_institucional()
                
                # SECCIÓN 3: CARACTERIZACIÓN
                self._agregar_modalidades_duracion()
                self._agregar_informacion_matriculas()
                self._agregar_cobertura_geografica_detallada()
                
                # SECCIÓN 4: MERCADO
                self._agregar_estado_programas()
                self._agregar_analisis_competencia()
                self._agregar_demanda_mercado()
                self._agregar_tendencias_mercado()
            
            # Análisis de agentes
            self._agregar_analisis_denominacion()
            self._agregar_analisis_tendencias()
            
            # Finales
            self._agregar_graficas()
            self._agregar_conclusiones()
            self._agregar_recomendaciones()
            
            Path(filepath).parent.mkdir(parents=True, exist_ok=True)
            self.prs.save(filepath)
            print(f"✅ Presentación guardada en: {filepath}")
        
        except Exception as e:
            print(f"❌ Error generando presentación: {e}")
            import traceback
            traceback.print_exc()
    
    # ========== DIAPOSITIVAS BÁSICAS ==========
    
    def _agregar_portada(self) -> None:
        """Agrega diapositiva de portada"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.color_primario
        
        titulo_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        titulo_frame = titulo_box.text_frame
        titulo_frame.word_wrap = True
        p = titulo_frame.paragraphs[0]
        p.text = "Análisis de Oportunidad"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self.color_blanco
        p.alignment = PP_ALIGN.CENTER
        
        subtit_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        subtit_frame = subtit_box.text_frame
        p = subtit_frame.paragraphs[0]
        p.text = "Programas Académicos SNIES - Análisis Completo"
        p.font.size = Pt(28)
        p.font.color.rgb = self.color_blanco
        p.alignment = PP_ALIGN.CENTER
        
        prog_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
        prog_frame = prog_box.text_frame
        p = prog_frame.paragraphs[0]
        programa = self.resultados.get('programa', 'Programa Académico')
        p.text = f"Programa: {programa}"
        p.font.size = Pt(18)
        p.font.color.rgb = self.color_blanco
        p.alignment = PP_ALIGN.CENTER
        
        fecha_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        fecha_frame = fecha_box.text_frame
        p = fecha_frame.paragraphs[0]
        timestamp = self.resultados.get('timestamp', 'Noviembre 2025')
        p.text = timestamp
        p.font.size = Pt(14)
        p.font.color.rgb = self.color_blanco
        p.alignment = PP_ALIGN.CENTER
    
    def _agregar_tabla_contenidos(self) -> None:
        """Agrega tabla de contenidos completa"""
        slide = self._crear_slide_titulo("Contenido")
        
        contenidos = [
            "1. Resumen Ejecutivo",
            "2. Contexto del Mercado Académico",
            "3. Análisis de Programas (denominaciones)",
            "4. Instituciones Oferentes (universidades)",
            "5. Tecnológicas e Instituciones Técnicas",
            "6. Distribución Institucional",
            "7. Modalidades y Estructura Académica",
            "8. Información de Matrículas",
            "9. Cobertura Geográfica Detallada",
            "10. Análisis Competitivo",
            "11. Demanda de Mercado",
            "12. Análisis de Denominación (IA)",
            "13. Tendencias y Oportunidades",
            "14. Conclusiones y Recomendaciones"
        ]
        
        left = Inches(1.5)
        top = Inches(1.5)
        
        for i, contenido in enumerate(contenidos):
            text_box = slide.shapes.add_textbox(left, top + Inches(i * 0.38), Inches(7), Inches(0.35))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.paragraphs[0]
            p.text = contenido
            p.font.size = Pt(11)
            p.font.color.rgb = self.color_texto
    
    def _agregar_resumen_ejecutivo(self) -> None:
        """Agrega resumen ejecutivo"""
        slide = self._crear_slide_titulo("Resumen Ejecutivo")
        
        sintesis = self.resultados.get('sintesis', {})
        resumen = sintesis.get('resumen_ejecutivo', 'No disponible')
        
        # Estadísticas rápidas
        left = Inches(0.8)
        top = Inches(1.5)
        
        stats = [
            (f"Programas Equivalentes", self.datos_enriquecidos.get('cantidad_equivalentes', 'N/A')),
            (f"Instituciones Oferentes", self.datos_enriquecidos.get('instituciones', {}).get('total', 'N/A')),
            (f"Departamentos Cubiertos", self.datos_enriquecidos.get('cobertura_geografica', {}).get('total_departamentos', 'N/A')),
        ]
        
        for i, (label, valor) in enumerate(stats):
            label_box = slide.shapes.add_textbox(left, top + Inches(i * 0.5), Inches(3), Inches(0.4))
            label_frame = label_box.text_frame
            p = label_frame.paragraphs[0]
            p.text = label
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.color_primario
            
            valor_box = slide.shapes.add_textbox(left + Inches(3.5), top + Inches(i * 0.5), Inches(2), Inches(0.4))
            valor_frame = valor_box.text_frame
            p = valor_frame.paragraphs[0]
            p.text = str(valor)
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.color_acento
        
        # Resumen
        text_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(8.4), Inches(3.8))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        lineas = resumen.split('\n')[:10]
        for i, linea in enumerate(lineas):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = linea.strip()
            p.font.size = Pt(10)
            p.font.color.rgb = self.color_texto
            p.space_after = Pt(4)
    
    # ========== SECCIÓN 1: PROGRAMAS DETALLADO ==========
    
    def _agregar_programas_equivalentes_completo(self) -> None:
        """Agrega TODOS los programas equivalentes en detalle"""
        programas = self.datos_enriquecidos.get('programas_equivalentes', [])
        cantidad = len(programas)
        
        # Si hay muchos programas, crear múltiples diapositivas
        programas_por_slide = 20
        num_slides = (cantidad + programas_por_slide - 1) // programas_por_slide
        
        for slide_num in range(num_slides):
            titulo = f"Programas Equivalentes ({slide_num + 1}/{num_slides})"
            slide = self._crear_slide_titulo(titulo)
            
            inicio = slide_num * programas_por_slide
            fin = min((slide_num + 1) * programas_por_slide, cantidad)
            programas_slice = programas[inicio:fin]
            
            # Tabla con programas
            rows = len(programas_slice) + 1
            cols = 3
            left = Inches(0.5)
            top = Inches(1.3)
            width = Inches(9)
            height = Inches(0.35 * rows)
            
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Encabezados
            for col_idx, titulo_col in enumerate(['No.', 'Denominación', 'Tipo']):
                celda = table_shape.cell(0, col_idx)
                celda.text = titulo_col
                celda.fill.solid()
                celda.fill.fore_color.rgb = self.color_secundario
                para = celda.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.color.rgb = self.color_blanco
                para.font.size = Pt(8)
            
            # Datos
            for idx, prog in enumerate(programas_slice, 1):
                celda = table_shape.cell(idx, 0)
                celda.text = str(inicio + idx)
                celda.text_frame.paragraphs[0].font.size = Pt(8)
                
                celda = table_shape.cell(idx, 1)
                celda.text = str(prog)
                celda.text_frame.paragraphs[0].font.size = Pt(8)
                celda.text_frame.word_wrap = True
                
                # Determinar tipo
                tipo_prog = self._clasificar_tipo_programa(prog)
                celda = table_shape.cell(idx, 2)
                celda.text = tipo_prog
                celda.text_frame.paragraphs[0].font.size = Pt(8)
    
    def _clasificar_tipo_programa(self, programa: str) -> str:
        """Clasifica el tipo de programa"""
        prog_lower = str(programa).lower()
        if 'doctorado' in prog_lower:
            return 'Doctorado'
        elif 'maestria' in prog_lower or 'maestría' in prog_lower:
            return 'Maestría'
        elif 'especializ' in prog_lower:
            return 'Especialización'
        elif 'ingenieria' in prog_lower or 'ingeniería' in prog_lower:
            return 'Ingeniería'
        elif 'tecnologo' in prog_lower or 'tecnólogo' in prog_lower:
            return 'Tecnólogo'
        else:
            return 'Otro'
    
    def _agregar_denominaciones_analisis(self) -> None:
        """Análisis detallado de denominaciones"""
        slide = self._crear_slide_titulo("Análisis de Denominaciones")
        
        programas = self.datos_enriquecidos.get('programas_equivalentes', [])
        
        # Estadísticas
        left = Inches(0.8)
        top = Inches(1.5)
        
        # Contar tipos
        tipos = {}
        for prog in programas:
            tipo = self._clasificar_tipo_programa(prog)
            tipos[tipo] = tipos.get(tipo, 0) + 1
        
        stats_box = slide.shapes.add_textbox(left, top, Inches(8.4), Inches(0.6))
        stats_frame = stats_box.text_frame
        stats_frame.word_wrap = True
        p = stats_frame.paragraphs[0]
        p.text = f"Total de denominaciones encontradas: {len(programas)}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.color_primario
        
        # Tabla de tipos
        top = Inches(2.3)
        rows = len(tipos) + 1
        cols = 2
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, Inches(4), Inches(0.4 * rows)).table
        
        # Encabezados
        for col_idx, titulo in enumerate(['Tipo de Programa', 'Cantidad']):
            celda = table_shape.cell(0, col_idx)
            celda.text = titulo
            celda.fill.solid()
            celda.fill.fore_color.rgb = self.color_primario
            para = celda.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.color.rgb = self.color_blanco
            para.font.size = Pt(9)
        
        # Datos
        for row_idx, (tipo, cantidad) in enumerate(sorted(tipos.items(), key=lambda x: x[1], reverse=True), 1):
            celda = table_shape.cell(row_idx, 0)
            celda.text = tipo
            celda.text_frame.paragraphs[0].font.size = Pt(9)
            
            celda = table_shape.cell(row_idx, 1)
            celda.text = str(cantidad)
            celda.text_frame.paragraphs[0].font.size = Pt(9)
            celda.fill.solid()
            celda.fill.fore_color.rgb = RGBColor(230, 240, 250)
        
        # Palabras más frecuentes
        palabras_box = slide.shapes.add_textbox(Inches(5.5), Inches(2.3), Inches(3.9), Inches(4))
        palabras_frame = palabras_box.text_frame
        palabras_frame.word_wrap = True
        
        p = palabras_frame.paragraphs[0]
        p.text = "Palabras Clave:"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.color_primario
        
        # Extraer palabras
        todas_palabras = []
        for prog in programas:
            todas_palabras.extend(str(prog).lower().split())
        
        from collections import Counter
        frecuencias = Counter(todas_palabras)
        
        for palabra, freq in frecuencias.most_common(10):
            if len(palabra) > 2:
                p = palabras_frame.add_paragraph()
                p.text = f"• {palabra} ({freq})"
                p.font.size = Pt(9)
    
    # ========== SECCIÓN 2: INSTITUCIONES DETALLADO ==========
    
    def _agregar_instituciones_completo(self) -> None:
        """Agrega información COMPLETA de todas las instituciones"""
        slide = self._crear_slide_titulo("Instituciones Oferentes - Resumen General")
        
        instituciones = self.datos_enriquecidos.get('instituciones', {})
        inst_lista = instituciones.get('lista', [])
        
        # Estadísticas generales
        left = Inches(0.8)
        top = Inches(1.5)
        
        stats = [
            (f"Total de Instituciones", instituciones.get('total', 0)),
            (f"Acreditadas Alta Calidad", instituciones.get('acreditadas_alta_calidad', 0)),
            (f"Universidades", instituciones.get('por_tipo', {}).get('Universidad', 0)),
            (f"Tecnológicas", instituciones.get('por_tipo', {}).get('Tecnologica', 0)),
            (f"Instituciones Técnicas", instituciones.get('por_tipo', {}).get('Institucion Tecnica', 0)),
        ]
        
        for i, (label, valor) in enumerate(stats):
            label_box = slide.shapes.add_textbox(left, top + Inches(i * 0.5), Inches(3.5), Inches(0.4))
            label_frame = label_box.text_frame
            p = label_frame.paragraphs[0]
            p.text = label
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.color_primario
            
            valor_box = slide.shapes.add_textbox(left + Inches(3.8), top + Inches(i * 0.5), Inches(1.6), Inches(0.4))
            valor_frame = valor_box.text_frame
            p = valor_frame.paragraphs[0]
            p.text = str(valor)
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.color_acento
        
        # Distribución por sector
        sector = instituciones.get('por_sector', {})
        if sector:
            sector_box = slide.shapes.add_textbox(left + Inches(5.5), top, Inches(3.9), Inches(3))
            sector_frame = sector_box.text_frame
            sector_frame.word_wrap = True
            
            p = sector_frame.paragraphs[0]
            p.text = "Distribución por Sector:"
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = self.color_primario
            
            for sec, count in sector.items():
                p = sector_frame.add_paragraph()
                porcentaje = (count / instituciones.get('total', 1) * 100)
                p.text = f"• {sec}: {count} ({porcentaje:.1f}%)"
                p.font.size = Pt(9)
    
    def _agregar_universidades_oferentes(self) -> None:
        """Agrega TODAS las UNIVERSIDADES que ofrecen el programa"""
        instituciones = self.datos_enriquecidos.get('instituciones', {})
        inst_lista = [i for i in instituciones.get('lista', []) if i.get('tipo') == 'Universidad']
        
        if not inst_lista:
            return
        
        # Dividir en múltiples slides
        inst_por_slide = 15
        num_slides = (len(inst_lista) + inst_por_slide - 1) // inst_por_slide
        
        for slide_num in range(num_slides):
            titulo = f"Universidades Oferentes ({slide_num + 1}/{num_slides})"
            slide = self._crear_slide_titulo(titulo)
            
            inicio = slide_num * inst_por_slide
            fin = min((slide_num + 1) * inst_por_slide, len(inst_lista))
            inst_slice = inst_lista[inicio:fin]
            
            # Tabla detallada
            rows = len(inst_slice) + 1
            cols = 4
            left = Inches(0.4)
            top = Inches(1.3)
            width = Inches(9.2)
            height = Inches(0.35 * rows)
            
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Encabezados
            for col_idx, titulo_col in enumerate(['Institución', 'Ciudad', 'Acreditación', 'Web']):
                celda = table_shape.cell(0, col_idx)
                celda.text = titulo_col
                celda.fill.solid()
                celda.fill.fore_color.rgb = self.color_primario
                para = celda.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.color.rgb = self.color_blanco
                para.font.size = Pt(8)
            
            # Datos
            for idx, inst in enumerate(inst_slice, 1):
                # Institución
                celda = table_shape.cell(idx, 0)
                celda.text = inst.get('nombre', 'N/A')[:35]
                celda.text_frame.paragraphs[0].font.size = Pt(7)
                celda.text_frame.word_wrap = True
                
                # Ciudad
                celda = table_shape.cell(idx, 1)
                celda.text = inst.get('municipio', 'N/A')[:15]
                celda.text_frame.paragraphs[0].font.size = Pt(7)
                
                # Acreditación
                celda = table_shape.cell(idx, 2)
                acred = inst.get('acreditacion_alta_calidad', 'No')
                celda.text = "Sí" if acred == 'Si' else "No"
                celda.text_frame.paragraphs[0].font.size = Pt(7)
                if acred == 'Si':
                    celda.fill.solid()
                    celda.fill.fore_color.rgb = RGBColor(200, 240, 200)
                
                # Web (abreviado)
                celda = table_shape.cell(idx, 3)
                web = inst.get('web', 'N/A')
                celda.text = web[:20] if web else 'N/A'
                celda.text_frame.paragraphs[0].font.size = Pt(7)
    
    def _agregar_tecnologicas_oferentes(self) -> None:
        """Agrega TODAS las TECNOLÓGICAS que ofrecen el programa"""
        instituciones = self.datos_enriquecidos.get('instituciones', {})
        inst_lista = [i for i in instituciones.get('lista', []) if i.get('tipo') == 'Tecnologica']
        
        if not inst_lista:
            return
        
        inst_por_slide = 20
        num_slides = (len(inst_lista) + inst_por_slide - 1) // inst_por_slide
        
        for slide_num in range(num_slides):
            titulo = f"Tecnológicas Oferentes ({slide_num + 1}/{num_slides})"
            slide = self._crear_slide_titulo(titulo)
            
            inicio = slide_num * inst_por_slide
            fin = min((slide_num + 1) * inst_por_slide, len(inst_lista))
            inst_slice = inst_lista[inicio:fin]
            
            # Tabla
            rows = len(inst_slice) + 1
            cols = 3
            left = Inches(0.5)
            top = Inches(1.3)
            width = Inches(9)
            height = Inches(0.35 * rows)
            
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Encabezados
            for col_idx, titulo_col in enumerate(['Institución', 'Ciudad', 'Naturaleza']):
                celda = table_shape.cell(0, col_idx)
                celda.text = titulo_col
                celda.fill.solid()
                celda.fill.fore_color.rgb = self.color_secundario
                para = celda.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.color.rgb = self.color_blanco
                para.font.size = Pt(8)
            
            # Datos
            for idx, inst in enumerate(inst_slice, 1):
                celda = table_shape.cell(idx, 0)
                celda.text = inst.get('nombre', 'N/A')[:40]
                celda.text_frame.paragraphs[0].font.size = Pt(7)
                
                celda = table_shape.cell(idx, 1)
                celda.text = inst.get('municipio', 'N/A')[:15]
                celda.text_frame.paragraphs[0].font.size = Pt(7)
                
                celda = table_shape.cell(idx, 2)
                celda.text = inst.get('naturaleza', 'N/A')[:20]
                celda.text_frame.paragraphs[0].font.size = Pt(7)
    
    def _agregar_distribucion_institucional(self) -> None:
        """Distribución geográfica de instituciones"""
        slide = self._crear_slide_titulo("Distribución Geográfica de Instituciones")
        
        instituciones = self.datos_enriquecidos.get('instituciones', {})
        depts = instituciones.get('por_departamento', {})
        
        if not depts:
            return
        
        # Tabla de departamentos
        rows = min(len(depts) + 1, 16)
        cols = 2
        left = Inches(1.5)
        top = Inches(1.5)
        width = Inches(7)
        height = Inches(0.35 * rows)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Encabezados
        for col_idx, titulo in enumerate(['Departamento', 'Instituciones']):
            celda = table_shape.cell(0, col_idx)
            celda.text = titulo
            celda.fill.solid()
            celda.fill.fore_color.rgb = self.color_primario
            para = celda.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.color.rgb = self.color_blanco
            para.font.size = Pt(9)
        
        # Datos ordenados
        for idx, (dept, count) in enumerate(sorted(depts.items(), key=lambda x: x[1], reverse=True)[:rows-1], 1):
            celda = table_shape.cell(idx, 0)
            celda.text = str(dept)
            celda.text_frame.paragraphs[0].font.size = Pt(8)
            
            celda = table_shape.cell(idx, 1)
            celda.text = str(count)
            celda.text_frame.paragraphs[0].font.size = Pt(8)
            celda.fill.solid()
            celda.fill.fore_color.rgb = RGBColor(230, 240, 250)
    
    # ========== SECCIÓN 3: CARACTERIZACIÓN ==========
    
    def _agregar_modalidades_duracion(self) -> None:
        """Información completa de modalidades y duración"""
        slide = self._crear_slide_titulo("Modalidades y Estructura Académica")
        
        modalidades = self.datos_enriquecidos.get('modalidades', {})
        duracion = self.datos_enriquecidos.get('duracion', {})
        
        # MODALIDADES
        mod_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4.2), Inches(0.4))
        mod_frame = mod_box.text_frame
        p = mod_frame.paragraphs[0]
        p.text = f"Modalidades ({len(modalidades.get('disponibles', []))} tipos):"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.color_primario
        
        top = Inches(2.0)
        for i, (mod, count) in enumerate(sorted(modalidades.get('distribucion', {}).items(), key=lambda x: x[1], reverse=True)):
            box = slide.shapes.add_textbox(Inches(1), top + Inches(i * 0.4), Inches(4), Inches(0.35))
            frame = box.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            total_prog = sum(modalidades.get('distribucion', {}).values())
            porcentaje = (count / total_prog * 100) if total_prog > 0 else 0
            p.text = f"• {mod}: {count} ({porcentaje:.1f}%)"
            p.font.size = Pt(9)
            p.font.color.rgb = self.color_texto
        
        # DURACIÓN
        dur_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.2), Inches(0.4))
        dur_frame = dur_box.text_frame
        p = dur_frame.paragraphs[0]
        p.text = "Estructura Académica:"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.color_primario
        
        dur_info_box = slide.shapes.add_textbox(Inches(5.2), Inches(2.0), Inches(4.2), Inches(4))
        dur_info_frame = dur_info_box.text_frame
        dur_info_frame.word_wrap = True
        
        periodos_str = ', '.join(str(p) for p in duracion.get('periodos_disponibles', [])[:15])
        p = dur_info_frame.paragraphs[0]
        p.text = f"Periodos:\n{periodos_str if periodos_str else 'N/A'}"
        p.font.size = Pt(9)
        
        creditos = duracion.get('creditos_disponibles', [])
        if creditos:
            p = dur_info_frame.add_paragraph()
            creditos_str = ', '.join(str(c) for c in creditos[:10])
            p.text = f"\nCréditos:\n{creditos_str}"
            p.font.size = Pt(9)
        
        periodic = duracion.get('periodicidad', [])
        if periodic:
            p = dur_info_frame.add_paragraph()
            p.text = f"\nPeriodicidad:\n{', '.join(periodic)}"
            p.font.size = Pt(9)
    
    def _agregar_informacion_matriculas(self) -> None:
        """Información detallada de matrículas"""
        slide = self._crear_slide_titulo("Análisis de Matrículas y Costos")
        
        matriculas = self.datos_enriquecidos.get('matriculas', {})
        
        # Crear tabla
        rows = 7
        cols = 2
        left = Inches(2)
        top = Inches(1.8)
        width = Inches(6)
        height = Inches(3.5)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Encabezados
        for col_idx, titulo in enumerate(['Métrica', 'Valor']):
            celda = table_shape.cell(0, col_idx)
            celda.text = titulo
            celda.fill.solid()
            celda.fill.fore_color.rgb = self.color_primario
            para = celda.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.color.rgb = self.color_blanco
            para.font.size = Pt(10)
        
        # Datos
        datos_matricula = [
            ('Mínima', f"${matriculas.get('minima', 0):,.0f}" if matriculas.get('minima') else 'N/A'),
            ('Máxima', f"${matriculas.get('maxima', 0):,.0f}" if matriculas.get('maxima') else 'N/A'),
            ('Promedio', f"${matriculas.get('promedio', 0):,.0f}" if matriculas.get('promedio') else 'N/A'),
            ('Mediana', f"${matriculas.get('mediana', 0):,.0f}" if matriculas.get('mediana') else 'N/A'),
            ('Desviación Estándar', f"${matriculas.get('desv_estandar', 0):,.0f}" if matriculas.get('desv_estandar') else 'N/A'),
            ('Registros con Matrícula', str(matriculas.get('registros_con_matricula', 0))),
        ]
        
        for row_idx, (metrica, valor) in enumerate(datos_matricula, 1):
            celda = table_shape.cell(row_idx, 0)
            celda.text = metrica
            celda.text_frame.paragraphs[0].font.size = Pt(10)
            celda.text_frame.paragraphs[0].font.bold = True
            
            celda = table_shape.cell(row_idx, 1)
            celda.text = valor
            celda.text_frame.paragraphs[0].font.size = Pt(10)
            celda.fill.solid()
            celda.fill.fore_color.rgb = RGBColor(230, 240, 250)
    
    def _agregar_cobertura_geografica_detallada(self) -> None:
        """Cobertura geográfica muy detallada"""
        slide = self._crear_slide_titulo("Cobertura Geográfica - Departamentos")
        
        cobertura = self.datos_enriquecidos.get('cobertura_geografica', {})
        depts = cobertura.get('departamentos', {})
        
        # Resumen
        resumen_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(0.35))
        resumen_frame = resumen_box.text_frame
        p = resumen_frame.paragraphs[0]
        p.text = f"Programa disponible en {cobertura.get('total_departamentos', 0)} departamentos y {cobertura.get('total_municipios', 0)} municipios"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.color_primario
        
        # Tabla de departamentos
        rows = min(len(depts) + 1, 18)
        cols = 2
        left = Inches(0.8)
        top = Inches(1.9)
        width = Inches(8.4)
        height = Inches(0.32 * rows)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Encabezados
        for col_idx, titulo in enumerate(['Departamento', 'Programas Equivalentes']):
            celda = table_shape.cell(0, col_idx)
            celda.text = titulo
            celda.fill.solid()
            celda.fill.fore_color.rgb = self.color_secundario
            para = celda.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.color.rgb = self.color_blanco
            para.font.size = Pt(9)
        
        # Datos ordenados
        for idx, (dept, count) in enumerate(sorted(depts.items(), key=lambda x: x[1], reverse=True)[:rows-1], 1):
            celda = table_shape.cell(idx, 0)
            celda.text = str(dept)[:35]
            celda.text_frame.paragraphs[0].font.size = Pt(8)
            celda.text_frame.word_wrap = True
            
            celda = table_shape.cell(idx, 1)
            celda.text = str(count)
            celda.text_frame.paragraphs[0].font.size = Pt(8)
            celda.fill.solid()
            celda.fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    # ========== SECCIÓN 4: MERCADO ==========
    
    def _agregar_contexto_mercado(self) -> None:
        """Contexto completo del mercado académico"""
        slide = self._crear_slide_titulo("Contexto del Mercado Académico")
        
        ctx = self.datos_enriquecidos.get('contexto_mercado', {})
        
        # Crear tabla con información de mercado
        rows = 9
        cols = 2
        left = Inches(0.8)
        top = Inches(1.4)
        width = Inches(8.4)
        height = Inches(4.8)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Encabezados
        for col_idx, titulo in enumerate(['Métrica', 'Valor']):
            celda = table_shape.cell(0, col_idx)
            celda.text = titulo
            celda.fill.solid()
            celda.fill.fore_color.rgb = self.color_primario
            para = celda.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.color.rgb = self.color_blanco
            para.font.size = Pt(9)
        
        # Datos del mercado
        datos_mercado = [
            ('Instituciones Oferentes', str(ctx.get('competencia', {}).get('total_instituciones', 'N/A'))),
            ('Universidades', str(ctx.get('competencia', {}).get('universidades', 'N/A'))),
            ('Tecnológicas', str(ctx.get('competencia', {}).get('tecnologicas', 'N/A'))),
            ('Instituciones Técnicas', str(ctx.get('competencia', {}).get('instituciones_tecnicas', 'N/A'))),
            ('Sector Privado (%)', f"{ctx.get('sector', {}).get('porcentaje_privado', 'N/A')}%"),
            ('Acreditadas Alta Calidad (%)', f"{ctx.get('acreditacion', {}).get('porcentaje_acreditacion', 'N/A')}%"),
            ('Nivel de Demanda', ctx.get('demanda', {}).get('nivel_demanda', 'N/A')),
            ('Nuevos Inscritos Recientes', str(ctx.get('demanda', {}).get('total_nuevos_reciente', 'N/A'))),
        ]
        
        for row_idx, (metrica, valor) in enumerate(datos_mercado, 1):
            celda = table_shape.cell(row_idx, 0)
            celda.text = metrica
            celda.text_frame.paragraphs[0].font.size = Pt(8)
            celda.text_frame.paragraphs[0].font.bold = True
            
            celda = table_shape.cell(row_idx, 1)
            celda.text = valor
            celda.text_frame.paragraphs[0].font.size = Pt(8)
            celda.fill.solid()
            celda.fill.fore_color.rgb = RGBColor(230, 240, 250)
    
    def _agregar_estado_programas(self) -> None:
        """Estado de programas en detalle"""
        slide = self._crear_slide_titulo("Estado de Programas")
        
        estado = self.datos_enriquecidos.get('estado', {})
        
        left = Inches(0.8)
        top = Inches(1.5)
        
        # Estado programa
        title_box = slide.shapes.add_textbox(left, top, Inches(4), Inches(0.35))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Estado del Programa:"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.color_primario
        
        top_est = Inches(1.95)
        for i, (est, count) in enumerate(estado.get('estado_programa', {}).items()):
            box = slide.shapes.add_textbox(left, top_est + Inches(i * 0.4), Inches(4), Inches(0.35))
            frame = box.text_frame
            p = frame.paragraphs[0]
            p.text = f"• {est}: {count}"
            p.font.size = Pt(10)
            p.font.color.rgb = self.color_texto
        
        # Reconocimiento
        left2 = Inches(5.2)
        title_box2 = slide.shapes.add_textbox(left2, top, Inches(4), Inches(0.35))
        title_frame2 = title_box2.text_frame
        p = title_frame2.paragraphs[0]
        p.text = "Reconocimiento:"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.color_primario
        
        top_rec = Inches(1.95)
        reconocimientos = estado.get('reconocimiento', [])
        for i, rec in enumerate(reconocimientos[:6]):
            box = slide.shapes.add_textbox(left2, top_rec + Inches(i * 0.35), Inches(4), Inches(0.3))
            frame = box.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            p.text = f"• {rec}"
            p.font.size = Pt(9)
            p.font.color.rgb = self.color_texto
    
    def _agregar_analisis_competencia(self) -> None:
        """Análisis competitivo detallado"""
        slide = self._crear_slide_titulo("Análisis de Competencia")
        
        ctx = self.datos_enriquecidos.get('contexto_mercado', {})
        
        info_items = []
        
        comp = ctx.get('competencia', {})
        if comp:
            info_items.append(('Universidades', f"{comp.get('universidades', 0)} de {comp.get('total_instituciones', 0)}"))
            info_items.append(('Tecnológicas', f"{comp.get('tecnologicas', 0)} de {comp.get('total_instituciones', 0)}"))
            info_items.append(('Instituciones Técnicas', f"{comp.get('instituciones_tecnicas', 0)} de {comp.get('total_instituciones', 0)}"))
        
        sector = ctx.get('sector', {})
        if sector:
            info_items.append(('Sector Privado', f"{sector.get('porcentaje_privado', 0)}%"))
            info_items.append(('Sector Oficial', f"{100 - sector.get('porcentaje_privado', 0)}%"))
        
        acred = ctx.get('acreditacion', {})
        if acred:
            info_items.append(('Acreditadas Alta Calidad', f"{acred.get('porcentaje_acreditacion', 0)}%"))
        
        concentra = ctx.get('concentracion_geografica', {})
        if concentra:
            es_conc = "Sí (concentrada)" if concentra.get('es_concentrado', False) else "No (dispersa)"
            info_items.append(('Concentración Geográfica', es_conc))
        
        demanda = ctx.get('demanda', {})
        if demanda:
            info_items.append(('Demanda', demanda.get('nivel_demanda', 'N/A')))
            info_items.append(('Nuevos Inscritos', f"{demanda.get('total_nuevos_reciente', 0):,}"))
        
        mercado = ctx.get('precios', {})
        if mercado:
            homogeneo = "Sí (homogéneo)" if mercado.get('mercado_homogeneo', False) else "No (diferenciado)"
            info_items.append(('Mercado de Precios', homogeneo))
        
        # Mostrar items en columnas
        left_col = Inches(0.8)
        right_col = Inches(5.5)
        top = Inches(1.7)
        
        for i, (label, valor) in enumerate(info_items):
            if i % 2 == 0:
                col = left_col
                row_offset = (i // 2) * 0.65
            else:
                col = right_col
                row_offset = ((i - 1) // 2) * 0.65
            
            label_box = slide.shapes.add_textbox(col, top + Inches(row_offset), Inches(3.5), Inches(0.55))
            label_frame = label_box.text_frame
            label_frame.word_wrap = True
            p = label_frame.paragraphs[0]
            p.text = label + ":"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.color_primario
            
            valor_box = slide.shapes.add_textbox(col + Inches(0), top + Inches(row_offset + 0.3), Inches(3.5), Inches(0.35))
            valor_frame = valor_box.text_frame
            valor_frame.word_wrap = True
            p = valor_frame.paragraphs[0]
            p.text = str(valor)
            p.font.size = Pt(10)
            p.font.color.rgb = self.color_texto
    
    def _agregar_demanda_mercado(self) -> None:
        """Análisis de demanda del mercado"""
        slide = self._crear_slide_titulo("Demanda de Mercado")
        
        ctx = self.datos_enriquecidos.get('contexto_mercado', {})
        demanda = ctx.get('demanda', {})
        evolucion = self.datos_enriquecidos.get('evolucion_temporal', {})
        
        left = Inches(0.8)
        top = Inches(1.5)
        
        demand_info = [
            ('Nivel de Demanda', demanda.get('nivel_demanda', 'N/A')),
            ('Nuevos Inscritos (período reciente)', f"{demanda.get('total_nuevos_reciente', 0):,}"),
            ('Total Matriculados (período reciente)', f"{demanda.get('total_matriculados_reciente', 0):,}"),
            ('Período de Análisis', evolucion.get('rango_temporal', 'N/A')),
        ]
        
        for i, (label, valor) in enumerate(demand_info):
            label_box = slide.shapes.add_textbox(left, top + Inches(i * 0.55), Inches(4), Inches(0.5))
            label_frame = label_box.text_frame
            label_frame.word_wrap = True
            p = label_frame.paragraphs[0]
            p.text = label + ":"
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.color_primario
            
            valor_box = slide.shapes.add_textbox(left + Inches(4.2), top + Inches(i * 0.55), Inches(4.2), Inches(0.5))
            valor_frame = valor_box.text_frame
            valor_frame.word_wrap = True
            p = valor_frame.paragraphs[0]
            p.text = str(valor)
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = self.color_acento
    
    def _agregar_tendencias_mercado(self) -> None:
        """Tendencias del mercado académico"""
        slide = self._crear_slide_titulo("Tendencias del Mercado Académico")
        
        ctx = self.datos_enriquecidos.get('contexto_mercado', {})
        tendencias = ctx.get('tendencias', {})
        
        left = Inches(0.8)
        top = Inches(1.5)
        
        p_box = slide.shapes.add_textbox(left, top, Inches(8.4), Inches(5.5))
        p_frame = p_box.text_frame
        p_frame.word_wrap = True
        
        p = p_frame.paragraphs[0]
        p.text = "Análisis de Tendencias:"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = self.color_primario
        
        if tendencias.get('programas_crecientes'):
            p = p_frame.add_paragraph()
            p.text = "✓ Programas con tendencia CRECIENTE en el mercado"
            p.font.size = Pt(11)
            p.font.color.rgb = self.color_exito
            p.level = 0
        
        if tendencias.get('modalidades_modernas'):
            p = p_frame.add_paragraph()
            p.text = "✓ Disponibilidad de MODALIDADES MODERNAS (virtual/mixta)"
            p.font.size = Pt(11)
            p.font.color.rgb = self.color_exito
            p.level = 0
        
        p = p_frame.add_paragraph()
        p.text = "\nOportunidades de Mercado:"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.color_primario
        
        p = p_frame.add_paragraph()
        p.text = "• Alta demanda de especialización en esta área"
        p.font.size = Pt(10)
        
        p = p_frame.add_paragraph()
        p.text = "• Mercado diversificado con múltiples oferentes"
        p.font.size = Pt(10)
        
        p = p_frame.add_paragraph()
        p.text = "• Opciones de modalidades variadas"
        p.font.size = Pt(10)
        
        p = p_frame.add_paragraph()
        p.text = "• Programas en múltiples ubicaciones geográficas"
        p.font.size = Pt(10)
    
    # ========== DIAPOSITIVAS DE AGENTES ==========
    
    def _agregar_analisis_denominacion(self) -> None:
        """Análisis de denominación de agentes"""
        slide = self._crear_slide_titulo("Análisis de Denominación (IA)")
        
        denom = self.resultados.get('denominacion', {})
        analisis_ia = denom.get('analisis_ia', {})
        
        left = Inches(0.8)
        top = Inches(1.5)
        
        prog_box = slide.shapes.add_textbox(left, top, Inches(8.4), Inches(0.45))
        prog_frame = prog_box.text_frame
        prog_frame.word_wrap = True
        p = prog_frame.paragraphs[0]
        denominacion = analisis_ia.get('denominacion_oficial', 'No disponible')
        p.text = f"Denominación oficial: {denominacion}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = self.color_primario
        
        clasificacion = analisis_ia.get('clasificacion', 'No disponible')
        clase_box = slide.shapes.add_textbox(left, top + Inches(0.55), Inches(8.4), Inches(0.4))
        clase_frame = clase_box.text_frame
        p = clase_frame.paragraphs[0]
        p.text = f"Clasificación: {clasificacion}"
        p.font.size = Pt(10)
        p.font.color.rgb = self.color_texto
        
        if 'posicionamiento_mercado' in analisis_ia:
            pos = analisis_ia['posicionamiento_mercado']
            pos_box = slide.shapes.add_textbox(left, top + Inches(1.05), Inches(8.4), Inches(0.4))
            pos_frame = pos_box.text_frame
            pos_frame.word_wrap = True
            p = pos_frame.paragraphs[0]
            p.text = f"Competencia: {pos.get('competencia_nivel', 'N/A')} | Saturación: {pos.get('saturacion', 'N/A')} | Tendencia: {pos.get('tendencia', 'N/A')}"
            p.font.size = Pt(9)
            p.font.color.rgb = self.color_texto
        
        hall_box = slide.shapes.add_textbox(left, top + Inches(1.55), Inches(8.4), Inches(4.8))
        hall_frame = hall_box.text_frame
        hall_frame.word_wrap = True
        
        p = hall_frame.paragraphs[0]
        p.text = "Hallazgos Principales:"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = self.color_primario
        
        for hallazgo in analisis_ia.get('hallazgos', [])[:6]:
            p = hall_frame.add_paragraph()
            p.text = f"• {hallazgo}"
            p.font.size = Pt(9)
            p.level = 0
    
    def _agregar_analisis_tendencias(self) -> None:
        """Análisis de tendencias de agentes"""
        slide = self._crear_slide_titulo("Tendencias y Oportunidades (IA)")
        
        tend = self.resultados.get('tendencias', {})
        analisis_ia = tend.get('analisis_ia', {})
        
        emergentes = analisis_ia.get('emergentes', [])
        decadentes = analisis_ia.get('decadentes', [])
        
        # Palabras emergentes
        emerg_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(4.2), Inches(0.35))
        emerg_frame = emerg_box.text_frame
        p = emerg_frame.paragraphs[0]
        p.text = f"Palabras Emergentes ({len(emergentes)}):"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = self.color_exito
        
        emerg_list = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(4.2), Inches(4.8))
        emerg_frame = emerg_list.text_frame
        emerg_frame.word_wrap = True
        
        for i, palabra in enumerate(emergentes[:12]):
            if i == 0:
                p = emerg_frame.paragraphs[0]
            else:
                p = emerg_frame.add_paragraph()
            p.text = f"• {palabra}"
            p.font.size = Pt(8)
        
        # Palabras decadentes
        decad_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.2), Inches(0.35))
        decad_frame = decad_box.text_frame
        p = decad_frame.paragraphs[0]
        p.text = f"Palabras en Declive ({len(decadentes)}):"
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = self.color_advertencia
        
        decad_list = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.2), Inches(4.8))
        decad_frame = decad_list.text_frame
        decad_frame.word_wrap = True
        
        for i, palabra in enumerate(decadentes[:12]):
            if i == 0:
                p = decad_frame.paragraphs[0]
            else:
                p = decad_frame.add_paragraph()
            p.text = f"• {palabra}"
            p.font.size = Pt(8)
    
    # ========== DIAPOSITIVAS FINALES ==========
    
    def _agregar_graficas(self) -> None:
        """Agrega gráficas"""
        print("  Buscando gráficas generadas...")
        
        grafica_files = []
        if os.path.exists(self.graficas_dir):
            for file in os.listdir(self.graficas_dir):
                if file.startswith('grafica_') and file.endswith('.png'):
                    grafica_files.append(os.path.join(self.graficas_dir, file))
        
        if not grafica_files:
            print("  ⚠️  No se encontraron gráficas")
            return
        
        grafica_files.sort()
        print(f"  ✅ Se encontraron {len(grafica_files)} gráficas")
        
        for i in range(0, min(len(grafica_files), 4), 2):
            slide = self._crear_slide_vacio()
            
            try:
                slide.shapes.add_picture(grafica_files[i], Inches(0.3), Inches(0.5), width=Inches(4.7))
            except Exception as e:
                print(f"  ⚠️  Error añadiendo gráfica {i}: {e}")
            
            if i + 1 < len(grafica_files):
                try:
                    slide.shapes.add_picture(grafica_files[i + 1], Inches(5.2), Inches(0.5), width=Inches(4.7))
                except Exception as e:
                    print(f"  ⚠️  Error añadiendo gráfica {i+1}: {e}")
    
    def _agregar_conclusiones(self) -> None:
        """Conclusiones"""
        slide = self._crear_slide_titulo("Conclusiones")
        
        sintesis = self.resultados.get('sintesis', {})
        hallazgos = sintesis.get('hallazgos_principales', [])
        
        hall_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        hall_frame = hall_box.text_frame
        hall_frame.word_wrap = True
        
        for i, hallazgo in enumerate(hallazgos[:8]):
            if i == 0:
                p = hall_frame.paragraphs[0]
            else:
                p = hall_frame.add_paragraph()
            p.text = f"• {hallazgo}"
            p.font.size = Pt(10)
            p.font.color.rgb = self.color_texto
            p.space_after = Pt(8)
    
    def _agregar_recomendaciones(self) -> None:
        """Recomendaciones"""
        slide = self._crear_slide_titulo("Recomendaciones")
        
        sintesis = self.resultados.get('sintesis', {})
        recomendaciones = sintesis.get('recomendaciones', [])
        
        rec_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
        rec_frame = rec_box.text_frame
        rec_frame.word_wrap = True
        
        for i, recom in enumerate(recomendaciones[:8]):
            if i == 0:
                p = rec_frame.paragraphs[0]
            else:
                p = rec_frame.add_paragraph()
            p.text = f"• {recom}"
            p.font.size = Pt(10)
            p.font.color.rgb = self.color_texto
            p.space_after = Pt(8)
    
    # ========== UTILIDADES ==========
    
    def _crear_slide_titulo(self, titulo: str):
        """Crea un slide con título"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.color_blanco
        
        linea = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
        linea.fill.solid()
        linea.fill.fore_color.rgb = self.color_primario
        linea.line.color.rgb = self.color_primario
        
        titulo_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(8.4), Inches(0.55))
        titulo_frame = titulo_box.text_frame
        titulo_frame.word_wrap = True
        p = titulo_frame.paragraphs[0]
        p.text = titulo
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = self.color_primario
        
        return slide
    
    def _crear_slide_vacio(self):
        """Crea un slide en blanco"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.color_blanco
        
        return slide